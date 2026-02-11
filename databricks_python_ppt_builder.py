"""
============================================================================
DATABRICKS PYTHON SCRIPT: POWERPOINT PRESENTATION BUILDER
============================================================================
Purpose: Read data from R Spark temp views and create PowerPoint presentation
Input: Spark temp views + PNG charts from /dbfs/FileStore/charts/
Output: PowerPoint presentation in /dbfs/FileStore/presentations/
============================================================================
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime
import os

print("=" * 80)
print("SPORTS VIEWING ANALYTICS - POWERPOINT PRESENTATION BUILDER")
print("=" * 80)
print()

# ============================================================================
# SECTION 1: LOAD DATA FROM R SPARK TEMP VIEWS
# ============================================================================
print("1. Loading data from R Spark temp views...")

try:
    summary_data = spark.sql("SELECT * FROM summary_data").toPandas()
    monthly_data = spark.sql("SELECT * FROM monthly_data").toPandas()
    sports_data = spark.sql("SELECT * FROM sports_data").toPandas()
    device_data = spark.sql("SELECT * FROM device_data").toPandas()
    
    # Convert summary to dictionary
    summary = summary_data.iloc[0].to_dict()
    
    print("   ✓ Summary data loaded")
    print("   ✓ Monthly data loaded:", len(monthly_data), "rows")
    print("   ✓ Sports data loaded:", len(sports_data), "rows")
    print("   ✓ Device data loaded:", len(device_data), "rows")
    print()
    
except Exception as e:
    print(f"   ✗ Error loading data: {e}")
    print("   Make sure to run the R script first!")
    raise

# Display sample data
print("Summary Data:")
print(summary)
print("\nSports Data:")
print(sports_data)
print("\nDevice Data:")
print(device_data)
print()

# ============================================================================
# SECTION 2: GENERATE INSIGHTS & RECOMMENDATIONS
# ============================================================================
print("2. Generating insights and recommendations...")

def generate_key_findings(summary, monthly_data, sports_data, device_data):
    """Generate key findings from the data"""
    findings = []
    
    # Overall engagement
    findings.append(
        f"Total viewing time reached {int(summary['TotalViewingMinutes']):,} minutes "
        f"across {int(summary['UniqueViewers']):,} unique viewers"
    )
    
    # Top competition
    top_comp = monthly_data.groupby('Competition')['ViewingMinutes'].sum().idxmax()
    top_comp_views = int(monthly_data.groupby('Competition')['ViewingMinutes'].sum().max())
    findings.append(
        f"{top_comp} led all competitions with {top_comp_views:,} total viewing minutes"
    )
    
    # Engagement rate
    avg_engagement = summary['AvgMinutesPerViewer']
    findings.append(
        f"Average viewer engagement of {avg_engagement:.1f} minutes per viewer "
        f"indicates strong content retention"
    )
    
    # Top sport
    top_sport = sports_data.nlargest(1, 'ViewingMinutes')['Sport'].values[0]
    findings.append(
        f"{top_sport} dominated sports viewership across all categories"
    )
    
    # Device preference
    top_device = device_data.nlargest(1, 'UniqueViewers')['Device'].values[0]
    findings.append(
        f"{top_device} was the preferred viewing device, capturing the largest audience share"
    )
    
    # Growth trend
    monthly_totals = monthly_data.groupby('YearMonth')['ViewingMinutes'].sum()
    first_half = monthly_totals.iloc[:6].mean()
    second_half = monthly_totals.iloc[6:].mean()
    growth = ((second_half - first_half) / first_half) * 100
    findings.append(
        f"Viewing minutes {'increased' if growth > 0 else 'decreased'} by "
        f"{abs(growth):.1f}% in the second half of the year"
    )
    
    return findings

def generate_recommendations(monthly_data, sports_data, device_data):
    """Generate strategic recommendations"""
    recommendations = []
    
    # Competition strategy
    top_3_comps = monthly_data.groupby('Competition')['UniqueViewers'].sum().nlargest(3).index.tolist()
    recommendations.append(
        f"Focus marketing efforts on top-performing competitions: {', '.join(top_3_comps)}"
    )
    
    # Device optimization
    mobile_users = device_data[device_data['Device'] == 'Mobile']['UniqueViewers'].values
    tv_users = device_data[device_data['Device'] == 'TV']['UniqueViewers'].values
    
    if len(mobile_users) > 0 and len(tv_users) > 0:
        if mobile_users[0] > tv_users[0]:
            recommendations.append(
                "Prioritize mobile app enhancements and responsive design given strong mobile adoption"
            )
        else:
            recommendations.append(
                "Optimize TV viewing experience with enhanced picture quality and interactive features"
            )
    
    # Content gaps
    low_engagement_sports = sports_data.nsmallest(2, 'MinutesPerViewer')['Sport'].tolist()
    if len(low_engagement_sports) >= 2:
        recommendations.append(
            f"Improve content quality and promotion for {' and '.join(low_engagement_sports)} "
            f"to boost engagement"
        )
    
    # Seasonal planning
    recommendations.append(
        "Develop targeted campaigns for Q4 when football competitions peak in viewership"
    )
    
    # Multi-platform strategy
    recommendations.append(
        "Implement seamless cross-device experiences to support viewers who switch between TV and mobile"
    )
    
    # Retention strategy
    recommendations.append(
        "Launch loyalty programs to convert casual viewers into regular subscribers"
    )
    
    return recommendations

# Generate insights
findings = generate_key_findings(summary, monthly_data, sports_data, device_data)
recommendations = generate_recommendations(monthly_data, sports_data, device_data)

print("   ✓ Generated", len(findings), "key findings")
print("   ✓ Generated", len(recommendations), "strategic recommendations")
print()

# ============================================================================
# SECTION 3: CREATE POWERPOINT PRESENTATION
# ============================================================================
print("3. Creating PowerPoint presentation...")

def create_presentation(summary, sports_data, device_data, findings, recommendations):
    """Create the PowerPoint presentation"""
    
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Define color scheme - Teal Trust theme
    colors = {
        'primary': RGBColor(2, 128, 144),      # #028090
        'secondary': RGBColor(0, 168, 150),    # #00A896
        'accent': RGBColor(2, 195, 154),       # #02C39A
        'dark': RGBColor(1, 79, 84),           # #014F54
        'light': RGBColor(232, 244, 245),      # #E8F4F5
        'white': RGBColor(255, 255, 255),
        'text': RGBColor(31, 41, 55),
        'text_light': RGBColor(107, 114, 128)
    }
    
    print("   Creating slides...")
    
    # ========================================================================
    # SLIDE 1: TITLE SLIDE
    # ========================================================================
    print("      - Slide 1: Title")
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['primary']
    
    # Main title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "SPORTS VIEWING"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle 1
    subtitle1_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(0.6))
    subtitle1_frame = subtitle1_box.text_frame
    subtitle1_frame.text = "ANALYTICS REPORT 2025"
    p = subtitle1_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle 2
    subtitle2_box = slide1.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(0.4))
    subtitle2_frame = subtitle2_box.text_frame
    subtitle2_frame.text = "Comprehensive Analysis of Viewing Trends, Engagement & Recommendations"
    p = subtitle2_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = colors['light']
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    p = footer_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = colors['light']
    p.alignment = PP_ALIGN.CENTER
    
    # ========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # ========================================================================
    print("      - Slide 2: Executive Summary")
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_shape = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "EXECUTIVE SUMMARY"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Metrics cards
    metrics = [
        ("📊", "Total Viewing Minutes", f"{int(summary['TotalViewingMinutes']):,}"),
        ("👥", "Unique Viewers", f"{int(summary['UniqueViewers']):,}"),
        ("⏱️", "Avg Minutes/Viewer", f"{summary['AvgMinutesPerViewer']:.1f}"),
        ("🎬", "Total Assets", f"{int(summary['TotalAssets']):,}")
    ]
    
    for idx, (icon, label, value) in enumerate(metrics):
        col = idx % 2
        row = idx // 2
        x = 0.5 + (col * 4.75)
        y = 1.2 + (row * 1.5)
        
        # Card background
        card = slide2.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.25), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = colors['light']
        card.line.fill.background()
        
        # Accent bar
        accent = slide2.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.15), Inches(1.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = colors['secondary']
        accent.line.fill.background()
        
        # Icon
        icon_box = slide2.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.15), Inches(0.6), Inches(0.6))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        p = icon_frame.paragraphs[0]
        p.font.size = Pt(36)
        
        # Value
        value_box = slide2.shapes.add_textbox(Inches(x + 1.1), Inches(y + 0.15), Inches(3), Inches(0.5))
        value_frame = value_box.text_frame
        value_frame.text = value
        p = value_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        
        # Label
        label_box = slide2.shapes.add_textbox(Inches(x + 1.1), Inches(y + 0.7), Inches(3), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = label
        p = label_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = colors['text_light']
    
    # ========================================================================
    # SLIDE 3: KEY FINDINGS
    # ========================================================================
    print("      - Slide 3: Key Findings")
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_shape = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "KEY FINDINGS"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    findings_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.5))
    findings_frame = findings_box.text_frame
    findings_frame.word_wrap = True
    
    for idx, finding in enumerate(findings):
        p = findings_frame.add_paragraph() if idx > 0 else findings_frame.paragraphs[0]
        p.text = finding
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(12)
    
    # ========================================================================
    # SLIDE 4: SPORTS PIE CHART
    # ========================================================================
    print("      - Slide 4: Sports Distribution Chart")
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_shape = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "VIEWERSHIP BY SPORT"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if os.path.exists("/dbfs/FileStore/charts/sports_pie.png"):
        slide4.shapes.add_picture("/dbfs/FileStore/charts/sports_pie.png", 
                                  Inches(1), Inches(1.2), width=Inches(8))
    else:
        print("         Warning: sports_pie.png not found")
    
    # ========================================================================
    # SLIDE 5: DEVICE LOLLIPOP CHART
    # ========================================================================
    print("      - Slide 5: Device Viewership Chart")
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_shape = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "VIEWERSHIP BY DEVICE"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if os.path.exists("/dbfs/FileStore/charts/device_lollipop.png"):
        slide5.shapes.add_picture("/dbfs/FileStore/charts/device_lollipop.png", 
                                  Inches(1), Inches(1.2), width=Inches(8))
    else:
        print("         Warning: device_lollipop.png not found")
    
    # ========================================================================
    # SLIDE 6: COMPETITION LINE CHART
    # ========================================================================
    print("      - Slide 6: Competition Trends Chart")
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_shape = slide6.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "COMPETITION TRENDS OVER TIME"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if os.path.exists("/dbfs/FileStore/charts/competition_line.png"):
        slide6.shapes.add_picture("/dbfs/FileStore/charts/competition_line.png", 
                                  Inches(0.5), Inches(1.2), width=Inches(9))
    else:
        print("         Warning: competition_line.png not found")
    
    # ========================================================================
    # SLIDE 7: TOP COMPETITIONS BAR CHART
    # ========================================================================
    print("      - Slide 7: Top Competitions Chart")
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_shape = slide7.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "TOP PERFORMING COMPETITIONS"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if os.path.exists("/dbfs/FileStore/charts/competition_bar.png"):
        slide7.shapes.add_picture("/dbfs/FileStore/charts/competition_bar.png", 
                                  Inches(1), Inches(1.2), width=Inches(8))
    else:
        print("         Warning: competition_bar.png not found")
    
    # ========================================================================
    # SLIDE 8: DATA TABLE
    # ========================================================================
    print("      - Slide 8: Data Table")
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_shape = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "DETAILED SPORTS METRICS"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Create table
    rows = len(sports_data) + 1
    cols = 4
    table_shape = slide8.shapes.add_table(rows, cols, Inches(0.8), Inches(1.2), 
                                          Inches(8.4), Inches(3.5))
    table = table_shape.table
    
    # Header row
    headers = ["Sport", "Viewing Minutes", "Unique Viewers", "Minutes/Viewer"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors['secondary']
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = colors['white']
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, (_, row) in enumerate(sports_data.iterrows(), start=1):
        table.cell(row_idx, 0).text = str(row['Sport'])
        table.cell(row_idx, 1).text = f"{int(row['ViewingMinutes']):,}"
        table.cell(row_idx, 2).text = f"{int(row['UniqueViewers']):,}"
        table.cell(row_idx, 3).text = f"{row['MinutesPerViewer']:.2f}"
        
        bg_color = colors['white'] if row_idx % 2 == 1 else colors['light']
        for col_idx in range(4):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER
    
    # ========================================================================
    # SLIDE 9: RECOMMENDATIONS
    # ========================================================================
    print("      - Slide 9: Strategic Recommendations")
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_shape = slide9.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = colors['primary']
    title_shape.line.fill.background()
    
    title_text = title_shape.text_frame
    title_text.text = "STRATEGIC RECOMMENDATIONS"
    p = title_text.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    title_text.margin_left = Inches(0.3)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    recs_box = slide9.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.5))
    recs_frame = recs_box.text_frame
    recs_frame.word_wrap = True
    
    for idx, rec in enumerate(recommendations):
        p = recs_frame.add_paragraph() if idx > 0 else recs_frame.paragraphs[0]
        p.text = f"{idx + 1}. {rec}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(12)
    
    # ========================================================================
    # SLIDE 10: CLOSING SLIDE
    # ========================================================================
    print("      - Slide 10: Closing")
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide10.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['primary']
    
    thank_you_box = slide10.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.8))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "THANK YOU"
    p = thank_you_frame.paragraphs[0]
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide10.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = colors['accent']
    p.alignment = PP_ALIGN.CENTER
    
    footer_box = slide10.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Sports Analytics Team | 2025"
    p = footer_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = colors['light']
    p.alignment = PP_ALIGN.CENTER
    
    # ========================================================================
    # SAVE PRESENTATION
    # ========================================================================
    print("   Saving presentation...")
    
    # Create output directory
    os.makedirs('/dbfs/FileStore/presentations', exist_ok=True)
    
    output_path = '/dbfs/FileStore/presentations/Sports_Viewing_Analytics_Report_2025.pptx'
    prs.save(output_path)
    
    return output_path

# Create the presentation
try:
    output_file = create_presentation(summary, sports_data, device_data, findings, recommendations)
    print("   ✓ Presentation saved successfully")
    print()
except Exception as e:
    print(f"   ✗ Error creating presentation: {e}")
    raise

# ============================================================================
# SECTION 4: COMPLETION SUMMARY
# ============================================================================
print("=" * 80)
print("POWERPOINT CREATION COMPLETE!")
print("=" * 80)
print()
print("✅ Presentation Details:")
print(f"   - Total Slides: 10")
print(f"   - File Location: {output_file}")
print(f"   - File Size: {os.path.getsize(output_file) / 1024:.1f} KB")
print()
print("📊 Content Summary:")
print(f"   - 1 Title slide")
print(f"   - 1 Executive summary with 4 key metrics")
print(f"   - 1 Key findings slide with {len(findings)} insights")
print(f"   - 4 Chart slides (sports, device, trends, top competitions)")
print(f"   - 1 Data table slide")
print(f"   - 1 Recommendations slide with {len(recommendations)} actions")
print(f"   - 1 Closing slide")
print()
print("📥 Download Instructions:")
print("   In Databricks, you can download the file from:")
print("   /FileStore/presentations/Sports_Viewing_Analytics_Report_2025.pptx")
print()
print("   Or use this link in your browser:")
print(f"   https://<your-databricks-instance>/files/presentations/Sports_Viewing_Analytics_Report_2025.pptx")
print()
print("=" * 80)

# Display download link (works in Databricks notebooks)
displayHTML(f"""
<div style="background-color: #E8F4F5; padding: 20px; border-radius: 10px; border: 2px solid #028090;">
    <h2 style="color: #028090; margin-top: 0;">✅ Presentation Ready!</h2>
    <p style="color: #1F2937; font-size: 16px;">
        Your comprehensive sports viewing analytics presentation has been created successfully.
    </p>
    <p style="color: #6B7280; font-size: 14px; margin-bottom: 20px;">
        Location: <code>/FileStore/presentations/Sports_Viewing_Analytics_Report_2025.pptx</code>
    </p>
    <a href="/files/presentations/Sports_Viewing_Analytics_Report_2025.pptx" download>
        <button style="background-color: #028090; color: white; padding: 12px 24px; 
                       border: none; border-radius: 5px; cursor: pointer; font-size: 16px;
                       font-weight: bold; box-shadow: 0 2px 4px rgba(0,0,0,0.2);">
            📥 Download PowerPoint Presentation
        </button>
    </a>
    <div style="margin-top: 20px; padding: 15px; background-color: white; border-radius: 5px;">
        <h4 style="color: #028090; margin-top: 0;">Presentation Contains:</h4>
        <ul style="color: #1F2937; line-height: 1.8;">
            <li><strong>{len(findings)}</strong> automated key findings</li>
            <li><strong>{len(recommendations)}</strong> strategic recommendations</li>
            <li><strong>4</strong> professional ggplot charts</li>
            <li><strong>1</strong> detailed data table</li>
            <li><strong>10</strong> total slides</li>
        </ul>
    </div>
</div>
""")
